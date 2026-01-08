import streamlit as st
import pandas as pd
import pdfplumber
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
import re
import json
import os
from datetime import datetime

# --- 1. 초기 세팅 및 데이터 저장 파일 설정 ---
st.set_page_config(page_title="Weekly Report Smart Converter", layout="wide")
HISTORY_FILE = "history_data.json"

def load_history_from_file():
    if os.path.exists(HISTORY_FILE):
        try:
            with open(HISTORY_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except:
            return []
    return []

def save_history_to_file(history):
    with open(HISTORY_FILE, "w", encoding="utf-8") as f:
        json.dump(history, f, ensure_ascii=False, indent=4)

if 'history' not in st.session_state:
    st.session_state['history'] = load_history_from_file()

# --- 2. 텍스트 간결화 및 중복 제거 함수 ---
def refine_text(text):
    if not text or str(text).lower() == 'nan' or text == "-": return "-"
    lines = str(text).split('\n')
    refined_lines = []
    seen = set()
    for line in lines:
        line = line.strip().replace('•', '').strip()
        if not line: continue
        # 비즈니스 문구 간결화
        line = re.sub(r' 진행 중(입니다)?', ' 진행', line)
        line = re.sub(r' 완료(하였습니다|했습니다)?', ' 완료', line)
        line = re.sub(r' 예정(입니다)?', ' 예정', line)
        line = line.replace(' 팔로업', ' F/U').replace('팔로우업', ' F/U')
        if line not in seen:
            refined_lines.append(f"• {line}")
            seen.add(line)
    return "\n".join(refined_lines) if refined_lines else "-"

# --- 3. 데이터 처리 로직 (엑셀/PDF 통합 개선) ---
def process_report_data(file):
    try:
        this_week_raw_list, next_week_raw_list = [], []
        
        if file.name.endswith('.pdf'):
            with pdfplumber.open(file) as pdf:
                for page in pdf.pages:
                    table = page.extract_table()
                    if not table: continue
                    df_tmp = pd.DataFrame(table)
                    # 헤더 찾기
                    h_idx = -1
                    for i, row in df_tmp.iterrows():
                        row_vals = [str(v) for v in row if v]
                        if any('프로젝트' in v or '팀원' in v for v in row_vals):
                            h_idx = i; break
                    if h_idx != -1:
                        for _, r in df_tmp.iloc[h_idx+1:].iterrows():
                            if len(r) >= 3 and r[1] and r[2]: this_week_raw_list.append([r[0], r[1], r[2]])
                            if len(r) >= 7 and r[5] and r[6]: next_week_raw_list.append([r[4], r[5], r[6]])
        else:
            # 엑셀 처리
            df_raw = pd.read_excel(file, sheet_name=0, header=None)
            h_idx = -1
            for i in range(len(df_raw)):
                row = [str(v).strip() for v in df_raw.iloc[i].values]
                if '프로젝트' in row or '팀원' in row:
                    h_idx = i; break
            
            if h_idx == -1:
                st.error("엑셀 파일에서 '프로젝트' 또는 '팀원' 헤더를 찾을 수 없습니다. 양식을 확인해 주세요.")
                return None
            
            data_df = df_raw.iloc[h_idx + 1:].copy()
            for _, r in data_df.iterrows():
                if len(r) >= 3 and pd.notna(r[1]) and str(r[1]).strip() != '':
                    this_week_raw_list.append([r[0], r[1], r[2]])
                if len(r) >= 7 and pd.notna(r[5]) and str(r[5]).strip() != '':
                    next_week_raw_list.append([r[4], r[5], r[6]])

        def summarize(rows):
            if not rows: return pd.DataFrame(columns=['프로젝트', '내용'])
            df = pd.DataFrame(rows, columns=['팀원', '프로젝트', '내용'])
            df['프로젝트'] = df['프로젝트'].astype(str).str.strip()
            df = df[~df['프로젝트'].str.contains('프로젝트|팀원|nan', case=False, na=False)]
            grouped = df.groupby('프로젝트')['내용'].apply(lambda x: refine_text("\n".join(map(str, x)))).reset_index()
            return grouped

        res_this = summarize(this_week_raw_list)
        res_next = summarize(next_week_raw_list)
        
        if res_this.empty and res_next.empty:
            st.warning("추출된 데이터가 없습니다. 파일 양식을 확인해 주세요.")
            return None

        merged = pd.merge(res_this, res_next, on='프로젝트', how='outer', suffixes=('_금', '_차')).fillna("-")
        merged.columns = ['프로젝트명', '이번 주 업무내용', '다음 주 업무내용']
        return merged.sort_values('프로젝트명')

    except Exception as e:
        st.error(f"파일 분석 중 오류: {e}")
        return None

# --- 4. PPT 생성 함수 ---
def create_split_pptx(df):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
    ROWS_PER_PAGE = 5 
    
    for i in range(0, len(df), ROWS_PER_PAGE):
        chunk = df.iloc[i : i + ROWS_PER_PAGE]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        p = title_box.text_frame.add_paragraph()
        p.text = f"서비스기획팀 주간업무보고 ({i//ROWS_PER_PAGE + 1})"
        p.font.bold, p.font.size = True, Pt(28)

        table = slide.shapes.add_table(len(chunk) + 1, 3, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.8)).table
        table.columns[0].width, table.columns[1].width, table.columns[2].width = Inches(2.3), Inches(5.0), Inches(5.0)

        headers = ["프로젝트명", "이번 주 업무내용", "다음 주 업무내용"]
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(44, 62, 80)
            p_head = cell.text_frame.paragraphs[0]
            p_head.font.color.rgb, p_head.font.bold, p_head.font.size = RGBColor(255, 255, 255), True, Pt(15)
            p_head.alignment = PP_ALIGN.CENTER

        for row_idx, (_, data) in enumerate(chunk.iterrows()):
            for col_idx in range(3):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(data.iloc[col_idx])
                for p_cell in cell.text_frame.paragraphs:
                    p_cell.font.size, p_cell.font.name = Pt(11), '맑은 고딕'
                    p_cell.alignment = PP_ALIGN.CENTER if col_idx == 0 else PP_ALIGN.LEFT

    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    return ppt_io.getvalue()

# --- 5. 사이드바 및 페이지 구성 ---
st.sidebar.title("📌 메뉴")
menu = st.sidebar.radio("이동할 페이지:", ["새 보고서 만들기", "변환 히스토리"])

if menu == "새 보고서 만들기":
    st.title("🚀 주간보고 스마트 PPT 변환기")
    file = st.file_uploader("Excel 또는 PDF 파일을 업로드하세요", type=["xlsx", "pdf"])

    if file:
        with st.spinner("데이터 분석 중..."):
            final_df = process_report_data(file)
            if final_df is not None:
                st.subheader("✅ 정제된 데이터 미리보기")
                st.dataframe(final_df, use_container_width=True)
                
                col1, col2 = st.columns(2)
                ppt_binary = create_split_pptx(final_df)
                
                with col1:
                    if st.button("💾 히스토리에 저장"):
                        history_item = {
                            "date": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                            "filename": file.name,
                            "data": final_df.to_dict('records')
                        }
                        st.session_state['history'].insert(0, history_item)
                        save_history_to_file(st.session_state['history'])
                        st.success("기록되었습니다!")
                
                with col2:
                    st.download_button(
                        label="📥 PPT 다운로드",
                        data=ppt_binary,
                        file_name=f"주간보고_정제본_{file.name.split('.')[0]}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )

elif menu == "변환 히스토리":
    st.title("📜 변환 히스토리")
    if not st.session_state['history']:
        st.info("저장된 이력이 없습니다.")
    else:
        # 개별 삭제 기능을 위해 반복문을 역순이 아닌 인덱스로 관리
        for idx, item in enumerate(st.session_state['history']):
            # 고유 키 생성을 위해 날짜와 인덱스 활용
            with st.expander(f"📅 {item['date']} - 📄 {item['filename']}"):
                hist_df = pd.DataFrame(item['data'])
                st.dataframe(hist_df, use_container_width=True)
                
                c1, c2 = st.columns([4, 1])
                with c1:
                    ppt_from_hist = create_split_pptx(hist_df)
                    st.download_button(
                        label=f"📥 PPT 다시 받기",
                        data=ppt_from_hist,
                        file_name=f"RE_{item['filename'].split('.')[0]}.pptx",
                        key=f"dl_{idx}"
                    )
                with c2:
                    # ★ 개별 삭제 버튼 추가
                    if st.button("❌ 기록 삭제", key=f"del_{idx}"):
                        st.session_state['history'].pop(idx)
                        save_history_to_file(st.session_state['history'])
                        st.rerun()
        
        st.divider()
        if st.sidebar.button("🗑️ 히스토리 전체 삭제"):
            if os.path.exists(HISTORY_FILE): os.remove(HISTORY_FILE)
            st.session_state['history'] = []
            st.rerun()